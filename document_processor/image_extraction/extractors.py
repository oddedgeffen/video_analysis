import os
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import tempfile
from PIL import Image
import io
import logging
import shutil

logger = logging.getLogger(__name__)

def extract_images_from_pdf(file_path, temp_dir):
    """Extract images from a PDF file."""
    images = []
    try:
        doc = fitz.open(file_path)
        for page_num in range(len(doc)):
            page = doc[page_num]
            image_list = page.get_images()
            
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    # Save in the provided temp directory
                    temp_path = os.path.join(temp_dir, f'pdf_page{page_num}_img{img_index}.png')
                    with open(temp_path, 'wb') as temp_file:
                        temp_file.write(image_bytes)
                        images.append(temp_path)
                        
                except Exception as e:
                    logger.error(f"Error extracting image {img_index} from PDF: {str(e)}")
                    continue
                    
        return images
    except Exception as e:
        logger.error(f"Error processing PDF file: {str(e)}")
        return []
    finally:
        if 'doc' in locals():
            doc.close()

def extract_images_from_docx(file_path, temp_dir):
    """Extract images from a DOCX file, including all types of headers and footers."""
    images = []
    try:
        doc = Document(file_path)
        
        # Extract images from the main document
        for idx, rel in enumerate(doc.part.rels.values()):
            if "image" in rel.reltype:
                try:
                    image_data = rel.target_part.blob
                    temp_path = os.path.join(temp_dir, f'docx_img{idx}.png')
                    with open(temp_path, 'wb') as temp_file:
                        temp_file.write(image_data)
                        images.append(temp_path)
                except Exception as e:
                    logger.error(f"Error extracting image from DOCX main document: {str(e)}")
                    continue
        
        # Extract images from headers and footers
        for section in doc.sections:
            # Process all types of headers (default, first page, even page)
            headers = [
                ('default', section.header),
                ('first_page', section.first_page_header),
                ('even_page', section.even_page_header)
            ]
            
            for header_type, header in headers:
                if header and header._element is not None:
                    for rel in header.part.rels.values():
                        if "image" in rel.reltype:
                            try:
                                image_data = rel.target_part.blob
                                temp_path = os.path.join(temp_dir, f'docx_{header_type}_header_img{len(images)}.png')
                                with open(temp_path, 'wb') as temp_file:
                                    temp_file.write(image_data)
                                    images.append(temp_path)
                            except Exception as e:
                                logger.error(f"Error extracting image from DOCX {header_type} header: {str(e)}")
                                continue
            
            # Process all types of footers (default, first page, even page)
            footers = [
                ('default', section.footer),
                ('first_page', section.first_page_footer),
                ('even_page', section.even_page_footer)
            ]
            
            for footer_type, footer in footers:
                if footer and footer._element is not None:
                    for rel in footer.part.rels.values():
                        if "image" in rel.reltype:
                            try:
                                image_data = rel.target_part.blob
                                temp_path = os.path.join(temp_dir, f'docx_{footer_type}_footer_img{len(images)}.png')
                                with open(temp_path, 'wb') as temp_file:
                                    temp_file.write(image_data)
                                    images.append(temp_path)
                            except Exception as e:
                                logger.error(f"Error extracting image from DOCX {footer_type} footer: {str(e)}")
                                continue
        
        return images
    except Exception as e:
        logger.error(f"Error processing DOCX file: {str(e)}")
        return []

def extract_images_from_pptx(file_path, temp_dir):
    """Extract images from a PPTX file."""
    images = []
    try:
        prs = Presentation(file_path)
        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "image"):
                    try:
                        image_bytes = shape.image.blob
                        temp_path = os.path.join(temp_dir, f'pptx_slide{slide_idx}_img{shape_idx}.png')
                        with open(temp_path, 'wb') as temp_file:
                            temp_file.write(image_bytes)
                            images.append(temp_path)
                    except Exception as e:
                        logger.error(f"Error extracting image from PPTX: {str(e)}")
                        continue
        return images
    except Exception as e:
        logger.error(f"Error processing PPTX file: {str(e)}")
        return []

def get_file_extension(file_path):
    """Get the lowercase file extension without the dot."""
    return os.path.splitext(file_path)[1].lower().lstrip('.')

def extract_images_from_document(file_path):
    """
    Extract images from a document based on its file type.
    Returns a tuple of (list of image paths, temp directory path).
    The caller is responsible for cleaning up the temp directory.
    """
    extension = get_file_extension(file_path)
    
    extractors = {
        'pdf': extract_images_from_pdf,
        'docx': extract_images_from_docx,
        'pptx': extract_images_from_pptx
    }
    
    if extension not in extractors:
        logger.error(f"Unsupported file type: {extension}")
        return [], None
    
    # Create a unique temporary directory for this extraction
    temp_dir = tempfile.mkdtemp()
    try:
        # Pass the temp directory to the appropriate extractor
        images = extractors[extension](file_path, temp_dir)
        return images, temp_dir
    except Exception as e:
        logger.error(f"Error extracting images from {extension} file: {str(e)}")
        # Clean up on error
        try:
            shutil.rmtree(temp_dir)
        except Exception as cleanup_error:
            logger.error(f"Error cleaning up temp directory: {str(cleanup_error)}")
        return [], None

def extract_images(document_path):
    images = []
    # Create a unique temporary subfolder
    with tempfile.TemporaryDirectory() as temp_dir:
        try:
            # ... existing code ...
                    # Save image in the unique temp subfolder with .png extension
                    temp_path = os.path.join(temp_dir, f'image_{len(images)}.png')
                    with open(temp_path, 'wb') as temp_file:
                        temp_file.write(image_bytes)
                        images.append(temp_path)
            # ... existing code ...
        except Exception as e:
            logger.error(f"Error extracting images: {str(e)}")
            # Clean up any remaining files
            for img_path in images:
                try:
                    os.unlink(img_path)
                except:
                    pass
            raise
    return images 